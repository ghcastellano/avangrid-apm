"""
PowerPoint Generator for Avangrid APM Platform
Generates a single PPTX with one slide per application, filled with data
from questionnaires, transcripts, and David's notes.
"""

import os
import re
import io
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt
from database import (
    get_session, close_session,
    Application, QuestionnaireAnswer, TranscriptAnswer, DavidNote, SynergyScore
)


TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "template.pptx")


# ============================================================
# Text transformation helpers (from generator_production.py)
# ============================================================

def clean_text(text):
    if text is None:
        return ""
    return str(text).strip()


def determine_satisfaction_emoji(answer):
    answer_lower = answer.lower()
    if any(w in answer_lower for w in ['high', 'good', 'positive', 'satisfied', 'excellent']):
        return "\U0001f60a"
    elif any(w in answer_lower for w in ['low', 'poor', 'negative', 'unsatisfied', 'bad']):
        return "\U0001f61e"
    return "\U0001f610"


def extract_yes_no(answer):
    answer_lower = answer.lower()
    if answer_lower.startswith('yes') or 'yes.' in answer_lower or 'yes,' in answer_lower:
        return "Yes"
    elif answer_lower.startswith('no') or 'no.' in answer_lower or 'no,' in answer_lower:
        return "No"
    return answer[:50]


def extract_criticality(answer):
    answer_lower = answer.lower()
    if 'critical' in answer_lower or 'high' in answer_lower or 'tier 1' in answer_lower or 'tier 2' in answer_lower:
        return "Yes (Business-Critical)"
    elif 'important' in answer_lower:
        return "Yes (Important)"
    elif 'supportive' in answer_lower:
        return "No (Supportive)"
    return answer[:50]


def extract_usage_trend(answer):
    answer_lower = answer.lower()
    if 'growing' in answer_lower or 'expansion' in answer_lower or 'increasing' in answer_lower:
        return "Growing"
    elif 'declining' in answer_lower or 'decreasing' in answer_lower:
        return "Declining"
    elif 'stable' in answer_lower:
        return "Stable"
    return answer[:50]


def extract_deployment(answer):
    answer_lower = answer.lower()
    if 'cloud' in answer_lower and 'on-prem' not in answer_lower and 'hybrid' not in answer_lower:
        return "Cloud"
    elif 'on-prem' in answer_lower and 'cloud' not in answer_lower:
        return "On-Premises"
    elif 'hybrid' in answer_lower:
        return "Hybrid"
    return answer[:50]


def extract_custom_or_market(answer):
    answer_lower = answer.lower()
    if 'custom' in answer_lower:
        return "Custom"
    elif 'market' in answer_lower or 'cots' in answer_lower or 'vendor' in answer_lower:
        return "Market"
    return answer[:50]


def extract_active_users(answer):
    match = re.search(r'(\d+)\s*active', answer.lower())
    if match:
        return match.group(1) + " active users"
    match2 = re.search(r'(\d[\d,]*)', answer)
    if match2:
        return match2.group(1) + " users"
    return answer[:50]


def extract_owned_by(answer):
    answer_lower = answer.lower()
    if 'business' in answer_lower and 'owned' in answer_lower:
        return "Business-Owned"
    elif 'it' in answer_lower and 'owned' in answer_lower:
        return "IT-Owned"
    return ""


# ============================================================
# Question-to-field keyword mapping
# ============================================================

def map_question_to_field(question_text, answer_text):
    """Map a question/answer pair to a PPT field using keyword matching.
    Returns (field_name, transformed_value) or (None, None) if no match."""
    q = question_text.lower()
    a = clean_text(answer_text)
    if not q or not a:
        return None, None

    if "name of the application" in q:
        return "NOME_DO_APP", a
    elif "utility domain" in q and "which" in q:
        return "UTILITY_DOMAIN", a
    elif "opco" in q and "which" in q:
        return "OPCOS", a
    elif "business unit" in q or ("business area" in q and "which" in q):
        return "BUSINESS_AREA", a
    elif "owned" in q and ("it-owned" in q or "business-owned" in q):
        return "OWNED_BY", extract_owned_by(a)
    elif "business-critical" in q or ("critical" in q and "application" in q):
        return "BUSINESS_CRITICALITY", extract_criticality(a)
    elif "usage" in q and ("growing" in q or "stable" in q or "declining" in q):
        return "USAGE", extract_usage_trend(a)
    elif "primary" in q and "purpose" in q:
        return "BUSINESS_PURPOSE", a
    elif "core functionalities" in q or ("functionalities" in q and "provide" in q):
        return "APPLICATION_FUNCTIONALITY", a
    elif "key business processes" in q and "support" in q:
        return "BUSINESS_PROCESS_COVERAGE", a
    elif "user roles" in q or "personas" in q:
        return "USER_PERSONAS", a
    elif "active users" in q and "how many" in q:
        return "NO_USERS", extract_active_users(a)
    elif "user satisfaction" in q or ("level of" in q and "satisfaction" in q):
        return "SATISFACTION", determine_satisfaction_emoji(a)
    elif "regulatory" in q and "compliance" in q and "support" in q:
        return "REGULATORY", extract_yes_no(a)
    elif "custom or market" in q or ("type of" in q and "application" in q):
        return "APP_CATEGORY", extract_custom_or_market(a)
    elif ("which systems" in q or "what systems" in q) and "integrate" in q:
        return "INTEGRATIONS", a
    elif "programming" in q:
        return "PROGRAMMING_LANGUAGE", a
    elif "deployed" in q or ("cloud" in q and "on" in q):
        return "DEPLOYMENT", extract_deployment(a)
    elif "platforms" in q or ("device" in q and "run" in q):
        return "COMPATIBILITY", a
    elif "monitored" in q or "apm" in q:
        return "MONITORED", extract_yes_no(a)
    elif "security" in q and ("risks" in q or "audit" in q or "gaps" in q):
        if a and a.lower() not in ('no', 'none'):
            return "SECURITY", "Yes"
        return "SECURITY", "No"
    elif "planned upgrades" in q or "migrations" in q or "replacements" in q:
        return "PLANNED_UPGRADES", a
    elif "user management" in q or "access control" in q or "provisioning" in q:
        return "USER_MANAGEMENT", a
    elif "bugs" in q or "incidents" in q or "issues" in q:
        return "BUGS_INCIDENTS", a
    elif "enhancements" in q or "improvements" in q:
        al = a.lower()
        if 'itnow' in al or 'email' in al:
            return "ENHANCEMENTS", "ITNow/email"
        elif 'vendor' in al or 'it' in al:
            return "ENHANCEMENTS", "IT, vendor, other"
        return "ENHANCEMENTS", a[:50]
    elif "business owner" in q or "business sponsor" in q:
        return "BUSINESS_OWNER", a
    elif "it owner" in q or "it manager" in q or "it lead" in q:
        return "IT_OWNER", a

    return None, None


# ============================================================
# Data extraction from database
# ============================================================

def extract_app_data_from_db(app_id, app_name, session=None):
    """Extract all data for an application from all database sources.
    Priority: DavidNote > TranscriptAnswer > QuestionnaireAnswer
    (later sources overwrite earlier ones, so highest-priority is last)

    If session is provided, uses it instead of creating a new one (avoids SQLite locking).
    """
    data = {
        "NOME_DO_APP": app_name,
        "UTILITY_DOMAIN": "", "OPCOS": "", "BUSINESS_AREA": "",
        "OWNED_BY": "", "BUSINESS_CRITICALITY": "", "USAGE": "",
        "BUSINESS_PURPOSE": "", "APPLICATION_FUNCTIONALITY": "",
        "BUSINESS_PROCESS_COVERAGE": "", "USER_PERSONAS": "",
        "NO_USERS": "", "BUSINESS_OWNER": "", "IT_OWNER": "",
        "SATISFACTION": "", "REGULATORY": "", "APP_CATEGORY": "",
        "INTEGRATIONS": "", "PROGRAMMING_LANGUAGE": "", "DEPLOYMENT": "",
        "COMPATIBILITY": "", "MONITORED": "", "SECURITY": "",
        "PLANNED_UPGRADES": "", "USER_MANAGEMENT": "",
        "BUGS_INCIDENTS": "", "ENHANCEMENTS": ""
    }

    own_session = session is None
    if own_session:
        session = get_session()
    try:
        # 1. Questionnaire answers (lowest priority - overwritten by others)
        qa_answers = session.query(QuestionnaireAnswer).filter_by(
            application_id=app_id
        ).all()
        for qa in qa_answers:
            field, value = map_question_to_field(qa.question_text or "", qa.answer_text or "")
            if field and value and field != "NOME_DO_APP":
                # Never overwrite NOME_DO_APP - we use the database app name
                data[field] = value

        # 2. Transcript answers (medium priority)
        ta_answers = session.query(TranscriptAnswer).filter_by(
            application_id=app_id
        ).all()
        for ta in ta_answers:
            field, value = map_question_to_field(ta.question_text or "", ta.answer_text or "")
            if field and value and field != "NOME_DO_APP":
                data[field] = value

        # 3. David's notes (highest priority - overwrites everything)
        david_notes = session.query(DavidNote).filter_by(
            application_id=app_id
        ).all()
        for dn in david_notes:
            field, value = map_question_to_field(dn.question_text or "", dn.answer_text or "")
            if field and value and field != "NOME_DO_APP":
                data[field] = value

        # Also try to extract owner info from David's notes that may have
        # different question patterns
        for dn in david_notes:
            q = (dn.question_text or "").lower()
            a = clean_text(dn.answer_text)
            if not a:
                continue
            if "owner" in q and "business" in q and not data["BUSINESS_OWNER"]:
                data["BUSINESS_OWNER"] = a
            elif "owner" in q and "it" in q and not data["IT_OWNER"]:
                data["IT_OWNER"] = a

    finally:
        if own_session:
            close_session(session)

    return data


# ============================================================
# Slide filling logic
# ============================================================

def auto_fit_text(text_frame, content, default_size=8, min_size=6, max_chars_per_size=None):
    """Set text with auto-fitting font size based on content length."""
    if max_chars_per_size is None:
        max_chars_per_size = {8: 150, 7: 200, 6: 280}

    font_size = default_size
    for size, max_chars in sorted(max_chars_per_size.items(), reverse=True):
        if len(content) <= max_chars:
            font_size = size
            break
    else:
        font_size = min_size
        content = content[:max_chars_per_size[min_size] - 3] + "..."

    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = content
    run.font.size = Pt(font_size)
    run.font.name = 'Calibri'
    run.font.bold = False
    text_frame.word_wrap = True


def set_text(text_frame, content, font_size=8, bold=False):
    """Set text with formatting, truncating if needed."""
    max_chars = 280
    if len(content) > max_chars:
        content = content[:max_chars - 3] + "..."
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = content
    run.font.size = Pt(font_size)
    run.font.name = 'Calibri'
    run.font.bold = bold
    text_frame.word_wrap = True


def fill_slide(slide, data_map):
    """Fill a single slide with application data (all 3 methods)."""

    # Step 0: Delete emoji images
    emoji_pictures = ['Picture 46', 'Picture 43', 'Picture 39']
    shapes_to_delete = []
    for shape in slide.shapes:
        if shape.name in emoji_pictures:
            shapes_to_delete.append(shape)
    for shape in shapes_to_delete:
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except Exception:
            pass

    # Step 1: Replace {{FIELD_NAME}} placeholders
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text.startswith('{{') and text.endswith('}}'):
                placeholder = text[2:-2]
                if placeholder == "NOME_DO_APP":
                    val = data_map.get(placeholder, "")
                    if val:
                        set_text(shape.text_frame, val, font_size=12, bold=True)
                    else:
                        shape.text_frame.clear()
                elif placeholder in data_map:
                    val = data_map.get(placeholder, "")
                    auto_fit_text(shape.text_frame, val)
        except Exception:
            pass

    # Step 2: Replace text-based placeholders
    placeholder_texts = {
        "<specify applications>": "INTEGRATIONS",
        "cloud or on prem": "DEPLOYMENT",
        "<per device type and os": "COMPATIBILITY",
        "custom or market": "APP_CATEGORY",
        "yes/no": None,  # handled by direct mapping
        "version, planned migrations": "PLANNED_UPGRADES",
    }
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip().lower()
            for placeholder_text, data_key in placeholder_texts.items():
                if placeholder_text in text and data_key and data_key in data_map:
                    val = data_map.get(data_key, "")
                    auto_fit_text(shape.text_frame, val)
                    break
        except Exception:
            pass

    # Step 3: Direct rectangle mapping
    direct_mappings = {
        "Rectangle 19": "APPLICATION_FUNCTIONALITY",
        "Rectangle 32": "BUSINESS_AREA",
        "Rectangle 36": "USER_PERSONAS",
        "Rectangle 142": "NO_USERS",
        "Rectangle 90": "REGULATORY",
        "Rectangle 64": "SATISFACTION",
        "Rectangle 189": "SECURITY",
        "Rectangle 174": "PROGRAMMING_LANGUAGE",
        "Rectangle 178": "PLANNED_UPGRADES",
        "Rectangle 182": "MONITORED",
    }
    for shape in slide.shapes:
        try:
            if shape.name in direct_mappings:
                data_key = direct_mappings[shape.name]
                if data_key in data_map and shape.has_text_frame:
                    val = data_map.get(data_key, "")
                    if data_key == "SATISFACTION":
                        set_text(shape.text_frame, val, font_size=16)
                    else:
                        auto_fit_text(shape.text_frame, val)
        except Exception:
            pass


# ============================================================
# Slide cloning with preserved relationship IDs
# ============================================================

def _copy_rels_with_rids(source_slide, new_slide, rids_needed):
    """Copy relationships from source to new slide, preserving exact rIds.
    This is critical: shape XML references rIds like rId5, rId6 etc.
    If the new slide has different rIds, images/embeds won't load."""
    import re as _re

    for rId, rel in source_slide.part.rels.items():
        if rId not in rids_needed:
            continue
        # Skip slideLayout rel (already set by add_slide)
        if 'slideLayout' in rel.reltype:
            continue
        try:
            target = rel.target_part
            if target is None:
                continue
            # Check if this rId already exists on new slide
            if rId in new_slide.part.rels:
                continue
            # Add relationship preserving the exact rId
            new_slide.part.rels._rels[rId] = rel
        except Exception:
            pass


# ============================================================
# Main generation function
# ============================================================

def generate_portfolio_pptx():
    """Generate a single PPTX with one slide per application.
    Returns bytes of the generated PPTX file."""
    import re as _re
    from lxml import etree

    session = get_session()
    try:
        apps = session.query(Application).order_by(Application.name).all()
        if not apps:
            return None

        # Extract data for all applications (pass session to avoid nested sessions)
        apps_data = []
        for app in apps:
            try:
                data = extract_app_data_from_db(app.id, app.name, session=session)
                apps_data.append(data)
            except Exception as e:
                print(f"[PPT_GENERATOR] Error extracting data for {app.name}: {e}")
                apps_data.append({"NOME_DO_APP": app.name})

        if not apps_data:
            return None

        # Load template
        prs = Presentation(TEMPLATE_PATH)

        # Save original template shapes XML for cloning
        template_slide = prs.slides[0]
        original_shapes = []
        for child in template_slide.shapes._spTree:
            tag_local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag_local in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
                original_shapes.append(deepcopy(child))

        # Discover which rIds are referenced by shape XML
        shapes_xml = etree.tostring(template_slide.shapes._spTree, encoding='unicode')
        referenced_rids = set(_re.findall(r'r:(?:embed|link|id)="(rId\d+)"', shapes_xml))

        # Fill first slide with first application
        try:
            fill_slide(template_slide, apps_data[0])
        except Exception as e:
            print(f"[PPT_GENERATOR] Error filling first slide ({apps_data[0].get('NOME_DO_APP', '?')}): {e}")

        # For remaining applications, clone template and fill
        for i in range(1, len(apps_data)):
            try:
                layout = template_slide.slide_layout
                new_slide = prs.slides.add_slide(layout)

                # Remove default shapes from layout
                for shape in list(new_slide.shapes):
                    new_slide.shapes._spTree.remove(shape.element)

                # Insert cloned original template shapes
                for shape_elem in original_shapes:
                    new_slide.shapes._spTree.append(deepcopy(shape_elem))

                # Copy ALL referenced relationships preserving exact rIds
                _copy_rels_with_rids(template_slide, new_slide, referenced_rids)

                # Fill with application data
                fill_slide(new_slide, apps_data[i])
            except Exception as e:
                print(f"[PPT_GENERATOR] Error creating slide {i} ({apps_data[i].get('NOME_DO_APP', '?')}): {e}")

        # Save to bytes buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    except Exception as e:
        print(f"[PPT_GENERATOR] Critical error generating PPTX: {e}")
        import traceback
        traceback.print_exc()
        return None
    finally:
        close_session(session)


def get_app_count():
    """Get number of applications that would be included in the PPT."""
    session = get_session()
    try:
        return session.query(Application).count()
    finally:
        close_session(session)
